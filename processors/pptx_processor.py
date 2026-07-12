"""
PowerPoint (.pptx) 处理器

提取所有幻灯片中的文本内容，保留标题和正文层次。
"""

import re
from pathlib import Path
from typing import Dict, Any, List

from .base_processor import BaseProcessor

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PptxProcessor(BaseProcessor):
    processor_name = "pptx"
    SUPPORTED_EXTENSIONS = {".pptx"}

    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == ".pptx"

    def extract_text(self, file_path: str) -> str:
        if not HAS_PPTX:
            raise ImportError(
                "缺少 python-pptx 库。请运行: pip install python-pptx"
            )

        prs = Presentation(file_path)
        text_parts: List[str] = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts: List[str] = []
            slide_texts.append(f"[幻灯片 {slide_idx}]")

            for shape in slide.shapes:
                shape_text = self._extract_shape_text(shape)
                if shape_text:
                    # 判断是否为标题
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        slide_texts.append(f"## {shape_text}")
                    else:
                        slide_texts.append(shape_text)

            if len(slide_texts) > 1:  # 有实际内容
                text_parts.append("\n".join(slide_texts))

        # 提取演讲者备注
        notes_texts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_texts.append(f"[幻灯片 {slide_idx} 备注]\n{notes}")

        if notes_texts:
            text_parts.append("\n\n---\n演讲者备注:\n" + "\n\n".join(notes_texts))

        return "\n\n".join(text_parts)

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        metadata = super().extract_metadata(file_path)
        try:
            prs = Presentation(file_path)
            metadata["slide_count"] = len(prs.slides)
            metadata["slide_width"] = prs.slide_width
            metadata["slide_height"] = prs.slide_height
        except Exception:
            pass
        return metadata

    # --- 内部方法 ---

    def _extract_shape_text(self, shape) -> str:
        """从 shape 中递归提取文本"""
        texts = []

        # 文本框
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    texts.append(para_text)

        # 表格
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append(" | ".join(row_texts))

        # 组合形状
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            if hasattr(shape, 'shapes'):
                for child in shape.shapes:
                    child_text = self._extract_shape_text(child)
                    if child_text:
                        texts.append(child_text)

        return "\n".join(texts)
